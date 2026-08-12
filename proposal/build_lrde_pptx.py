#!/usr/bin/env python3
"""Generate LRDE AI proposal PowerPoint from slide content."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt, Emu

# Palette
INK = RGBColor(0x0C, 0x16, 0x20)
MUTED = RGBColor(0x4A, 0x5A, 0x66)
PAPER = RGBColor(0xF4, 0xF1, 0xEB)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD2, 0xCD, 0xC4)
RADAR = RGBColor(0x0A, 0x6B, 0x52)
RADAR_DEEP = RGBColor(0x06, 0x45, 0x36)
STEEL = RGBColor(0x16, 0x32, 0x4F)
AMBER = RGBColor(0xC2, 0x41, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP_BG = RGBColor(0xE4, 0xF1, 0xEC)
DARK_BG = RGBColor(0x0C, 0x16, 0x20)
DARK_STEEL = RGBColor(0x16, 0x32, 0x4F)

W = Inches(13.333)
H = Inches(7.5)
MARGIN_X = Inches(0.7)
MARGIN_Y = Inches(0.45)


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE = 1
    fill_shape(shape, fill)
    return shape


def add_accent_bar(slide, left, top, height=Inches(0.9)):
    return add_rect(slide, left, top, Inches(0.08), height, RADAR)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_eyebrow(tf, text, color=RADAR):
    p = tf.paragraphs[0]
    p.text = text.upper()
    set_run(p.runs[0], 11, True, color)
    p.space_after = Pt(8)


def add_title(tf, text, color=STEEL, size=36):
    p = tf.add_paragraph()
    p.text = text
    set_run(p.runs[0], size, False, color, "Georgia")
    p.space_after = Pt(10)


def add_body(tf, text, color=MUTED, size=16):
    p = tf.add_paragraph()
    p.text = text
    set_run(p.runs[0], size, False, color)
    p.space_after = Pt(8)


def add_bullet(tf, text, color=MUTED, size=15, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    set_run(p.runs[0], size, False, color)
    p.space_after = Pt(4)


def card(slide, left, top, width, height, title, lines, accent=RADAR):
    add_rect(slide, left, top, width, height, PANEL)
    add_rect(slide, left, top, Inches(0.07), height, accent)
    box = add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.28), height - Inches(0.22))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_run(p.runs[0], 15, True, INK)
    p.space_after = Pt(6)
    for line in lines:
        add_body(tf, line, MUTED, 13)


def add_table(slide, left, top, width, rows, col_widths=None):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, Inches(0.38 * len(rows)))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            if r == 0:
                set_run(run, 11, True, INK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEB, 0xE6, 0xDE)
            else:
                set_run(run, 12, False, MUTED)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL
    return table_shape


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    prs = new_prs()
    content_w = W - 2 * MARGIN_X

    # 1 Title
    s = blank(prs)
    set_slide_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), W, H, DARK_BG)
    add_rect(s, Inches(9.5), Inches(0), Inches(3.833), H, RADAR)
    box = add_textbox(s, MARGIN_X, Inches(1.3), Inches(8.5), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    add_eyebrow(tf, "Proposal for LRDE · Bengaluru · Restricted", RGBColor(0xB8, 0xD9, 0xCC))
    p = tf.add_paragraph()
    p.text = "Altera"
    set_run(p.runs[0], 48, False, WHITE, "Georgia")
    p.space_after = Pt(6)
    add_title(tf, "AI Compute Platform for LRDE", WHITE, 40)
    add_body(
        tf,
        "Sovereign training hardware and Altera Agilex™ FPGA inference for radar signal intelligence, AESA processing aids, and real-time classification.",
        RGBColor(0xD5, 0xE4, 0xDC),
        16,
    )
    meta = add_textbox(s, MARGIN_X, Inches(6.4), content_w, Inches(0.7))
    mtf = meta.text_frame
    p = mtf.paragraphs[0]
    p.text = "Audience: LRDE leadership & radar AI teams    ·    Focus: Radar / AESA / EW-adjacent AI    ·    Version 1.0 · Aug 2026"
    set_run(p.runs[0], 12, False, RGBColor(0xB0, 0xC4, 0xBA))

    # 2 Agenda
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.5))
    tf = box.text_frame
    add_eyebrow(tf, "Agenda")
    add_title(tf, "What we will cover")
    add_body(tf, "A complete path from LRDE training data to fielded FPGA inference inside radar timelines.")
    items = [
        "LRDE mission fit & AI opportunities",
        "End-to-end architecture",
        "Training hardware (secure GPU lab)",
        "Inference on Altera Agilex FPGAs",
        "Software toolchain & MLOps",
        "LRDE use-case packages",
        "Security & sovereignty",
        "PoC roadmap & next steps",
    ]
    left_box = add_textbox(s, MARGIN_X, Inches(2.4), Inches(5.8), Inches(4.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:4], 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.text = f"{i}.  {item}"
        set_run(p.runs[0], 18, False, INK)
        p.space_after = Pt(14)
    right_box = add_textbox(s, Inches(7.2), Inches(2.4), Inches(5.4), Inches(4.2))
    tf = right_box.text_frame
    for i, item in enumerate(items[4:], 5):
        p = tf.paragraphs[0] if i == 5 else tf.add_paragraph()
        p.text = f"{i}.  {item}"
        set_run(p.runs[0], 18, False, INK)
        p.space_after = Pt(14)

    # 3 LRDE context
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    add_eyebrow(tf, "LRDE context")
    add_title(tf, "Radar AI where it matters")
    add_body(tf, "LRDE designs ground-based, shipborne, airborne, and space-related radar systems — AI must respect real-time DSP pipelines and SWaP.")
    card(s, MARGIN_X, Inches(2.5), Inches(3.8), Inches(3.6), "AESA & multifunction", [
        "Fire-control, surveillance, 4D radars.",
        "Multi-mode interleaved operation needs deterministic assistive AI.",
    ])
    card(s, Inches(4.75), Inches(2.5), Inches(3.8), Inches(3.6), "Signal intelligence", [
        "Pulse / emitter classification.",
        "LPI detection aids, clutter/anomaly separation on I/Q & spectrograms.",
    ])
    card(s, Inches(8.8), Inches(2.5), Inches(3.8), Inches(3.6), "Advanced research", [
        "Hypersonic cues, stealth-detection aids, FOPEN features.",
        "Netted-radar fusion — train secure, infer at the sensor.",
    ])

    # 4 Problem
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Problem")
    add_title(tf, "Gaps in today’s radar AI path")
    card(s, MARGIN_X, Inches(1.9), Inches(5.8), Inches(2.3), "Training", [
        "• Classified I/Q & track data cannot leave LRDE",
        "• Cloud GPUs unsuitable for air-gapped labs",
        "• Weak experiment / model lineage",
    ])
    card(s, MARGIN_X, Inches(4.4), Inches(5.8), Inches(2.3), "Inference", [
        "• GPUs struggle with radar SWaP & latency budgets",
        "• Need tight coupling to JESD / ADC / beam data",
        "• Long program life vs short GPU cycles",
    ])
    card(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.8), "Required outcome", [
        "A sovereign stack that lets LRDE:",
        "• Train on-prem on radar datasets",
        "• Compile models to FPGA AI IP",
        "• Deploy inside ERP / processor cabinets or payloads",
        "• Update bitstreams without board respins",
    ], accent=AMBER)

    # 5 Architecture
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Architecture")
    add_title(tf, "Train once · Infer on Agilex")
    flow = add_rect(s, MARGIN_X, Inches(1.9), content_w, Inches(3.5), INK)
    fbox = add_textbox(s, MARGIN_X + Inches(0.25), Inches(2.1), content_w - Inches(0.5), Inches(3.1))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    flow_text = (
        "LRDE SECURE LAB\n"
        "  Radar Data Vault → GPU Training Cluster → Model Registry\n"
        "                         │  ONNX / OpenVINO IR\n"
        "                         ▼\n"
        "  Quantize + Accuracy Gates → FPGA AI Suite (dla_compiler)\n"
        "          ┌──────────────┼──────────────┐\n"
        "   Hostless Agilex   SoC Agilex 5    PCIe Agilex 7\n"
        "   (ERP / payload)   (edge node)     (shelter / C2)"
    )
    p = ftf.paragraphs[0]
    p.text = flow_text
    set_run(p.runs[0], 14, False, RGBColor(0xD8, 0xEB, 0xE2), "Consolas")
    for i, (label, sub) in enumerate([("GPU", "Secure training"), ("OpenVINO", "IR + hetero runtime"), ("Agilex", "Field inference")]):
        left = MARGIN_X + i * Inches(4.1)
        add_rect(s, left, Inches(5.7), Inches(3.9), Inches(1.2), PANEL)
        t = add_textbox(s, left + Inches(0.2), Inches(5.85), Inches(3.5), Inches(0.9))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = label
        set_run(p.runs[0], 22, False, RADAR, "Georgia")
        add_body(tf, sub, MUTED, 13)

    # 6 Training HW
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    add_eyebrow(tf, "Training hardware")
    add_title(tf, "Air-gapped radar AI training lab")
    add_body(tf, "Sized for LRDE model development on spectrograms, RD maps, tracks, and multi-sensor fusion sets.")
    add_table(
        s,
        MARGIN_X,
        Inches(2.3),
        content_w,
        [
            ["Component", "Reference spec", "Role for LRDE"],
            ["GPU nodes", "8× datacenter GPUs/node, dual CPU, 1–2 TB RAM", "CNN / transformer training"],
            ["Fabric", "NVLink + InfiniBand NDR / RoCE", "Distributed training"],
            ["Storage", "Encrypted 0.5–2 PB object / parallel FS", "I/Q archives, checkpoints"],
            ["CPU ETL", "High-core labeling & feature nodes", "RD maps, CFAR labels, sim data"],
            ["Security", "Air gap, bastion MFA, optional diode", "Sovereign model export only"],
        ],
        [Inches(2.2), Inches(5.8), Inches(4.0)],
    )

    # 7 Training SW
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Training software")
    add_title(tf, "ML stack inside the wire")
    card(s, MARGIN_X, Inches(1.9), Inches(5.8), Inches(3.2), "Frameworks & ops", [
        "• PyTorch / TensorFlow (air-gapped registry)",
        "• DeepSpeed / FSDP for large models",
        "• Slurm or Kubernetes + Apptainer",
        "• MLflow / DVC for lineage",
    ])
    card(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(3.2), "Radar-specific data path", [
        "• I/Q → STFT / RD / spectrogram pipelines",
        "• Synthetic + trial data blending",
        "• Golden sets for accuracy gates",
        "• Signed ONNX / IR promotion",
    ])
    card(s, MARGIN_X, Inches(5.4), content_w, Inches(1.5), "Handoff to inference", [
        "Export ONNX → OpenVINO IR → FPGA AI Suite compile → signed bitstream / network-bin for LRDE platforms.",
    ])

    # 8 FPGA Inference
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    add_eyebrow(tf, "Inference hardware")
    add_title(tf, "Altera Agilex™ for radar timelines")
    add_body(tf, "FPGAs sit next to the exciter–receiver–processor chain — reconfigurable AI without breaking determinism.")
    add_table(
        s,
        MARGIN_X,
        Inches(2.3),
        content_w,
        [
            ["Device", "INT8 class*", "LRDE placement"],
            ["Agilex 3", "~2.6 TOPS", "Miniaturized / SoC radar sensors"],
            ["Agilex 5 E-Series", "~26 TOPS", "Compact processors, mobile radars"],
            ["Agilex 5 D-Series", "~152 TOPS", "Multi-channel classification"],
            ["Agilex 7 I/M-Series", "~89 TOPS class", "ERP co-process, PCIe, shelters"],
        ],
        [Inches(3.5), Inches(3.0), Inches(5.5)],
    )
    note = add_textbox(s, MARGIN_X, Inches(6.5), content_w, Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = "*Architecture- and IP-dependent; sized per model with FPGA AI Suite Architecture Optimizer."
    set_run(p.runs[0], 11, False, MUTED)

    # 9 Form factors
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Deployment")
    add_title(tf, "Form factors for LRDE systems")
    card(s, MARGIN_X, Inches(2.0), Inches(3.8), Inches(3.2), "Hostless / DDR-free", [
        "AI IP in fabric beside DSP.",
        "Lowest latency for pulse/window classifiers inside ERP.",
    ])
    card(s, Inches(4.75), Inches(2.0), Inches(3.8), Inches(3.2), "SoC (HPS)", [
        "On-chip Linux + AI IP.",
        "Ideal for transportable radars and edge fusion nodes.",
    ])
    card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(3.2), "PCIe / OFS card", [
        "Look-aside acceleration in operation shelters and lab validation racks.",
    ])
    chips = add_textbox(s, MARGIN_X, Inches(5.6), content_w, Inches(1.0))
    p = chips.text_frame.paragraphs[0]
    p.text = "JESD204 / ADC attach   ·   Aurora / Ethernet sensor nets   ·   Bitstream encryption   ·   Long defence lifecycle"
    set_run(p.runs[0], 14, True, RADAR_DEEP)

    # 10 Why FPGA
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Why Altera")
    add_title(tf, "FPGA vs GPU for LRDE inference")
    add_table(
        s,
        MARGIN_X,
        Inches(2.0),
        content_w,
        [
            ["Criterion", "GPU", "Altera FPGA"],
            ["Latency determinism", "Batch / driver variable", "Pipeline-friendly, stable"],
            ["RF / DSP coupling", "Needs host staging", "Direct fabric I/O"],
            ["Power in cabinet / payload", "High", "Lower W for fixed nets"],
            ["Model update", "Driver + app stack", "Bitstream / network-bin"],
            ["Program longevity", "Short SKU cycles", "10–15+ year class support"],
        ],
        [Inches(3.5), Inches(4.2), Inches(4.3)],
    )

    # 11 Software
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Software toolchain")
    add_title(tf, "From PyTorch to radar bitstream")
    add_rect(s, MARGIN_X, Inches(1.9), content_w, Inches(3.8), INK)
    fbox = add_textbox(s, MARGIN_X + Inches(0.3), Inches(2.15), content_w - Inches(0.6), Inches(3.4))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    flow_text = (
        "Train (PyTorch) → ONNX\n"
        "  → OpenVINO convert_model / ovc  →  IR (.xml/.bin)\n"
        "  → Quantize (INT8 / FP16) + golden accuracy gate\n"
        "  → FPGA AI Suite  dla_compiler --march <arch.arch>\n"
        "  → Architecture Optimizer (area ↔ FPS)\n"
        "  → Bit-accurate OpenVINO emulation\n"
        "  → Quartus Prime integrate AI IP → bitstream\n"
        "  → Runtime: FPGA  or  HETERO:FPGA,CPU  (or hostless)"
    )
    p = ftf.paragraphs[0]
    p.text = flow_text
    set_run(p.runs[0], 15, False, RGBColor(0xD8, 0xEB, 0xE2), "Consolas")
    note = add_textbox(s, MARGIN_X, Inches(6.1), content_w, Inches(0.6))
    p = note.text_frame.paragraphs[0]
    p.text = "Core tools: OpenVINO · FPGA AI Suite · Quartus Prime · optional Open FPGA Stack (OFS)."
    set_run(p.runs[0], 14, False, MUTED)

    # 12 Use cases
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "LRDE packages")
    add_title(tf, "Three PoC-ready use cases")
    card(s, MARGIN_X, Inches(2.0), Inches(3.8), Inches(4.0), "1. Emitter / pulse ID", [
        "Spectrogram CNN on Agilex 7 hostless or PCIe.",
        "Continuous windows, deterministic latency for EW-adjacent radar modes.",
    ])
    card(s, Inches(4.75), Inches(2.0), Inches(3.8), Inches(4.0), "2. Target / clutter assist", [
        "RD-map or tracklet classifier on Agilex 5.",
        "Aids discrimination for fighters, slow movers, false-alarm reduction.",
    ])
    card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(4.0), "3. Netted fusion node", [
        "Multi-radar feature fusion on Agilex 7 PCIe in shelter C2.",
        "Signed model hot-swap after trials.",
    ])

    # 13 Security
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Security")
    add_title(tf, "Sovereignty by design")
    card(s, MARGIN_X, Inches(2.0), Inches(5.8), Inches(2.0), "Data residency", [
        "I/Q, tracks, and labels never leave LRDE premises.",
    ])
    card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0), "Air gap", [
        "Training isolated; controlled diode export of IR/bitstreams only.",
    ])
    card(s, MARGIN_X, Inches(4.3), Inches(5.8), Inches(2.0), "Bitstream trust", [
        "Encryption, authentication, anti-tamper options on Agilex.",
    ])
    card(s, Inches(7.0), Inches(4.3), Inches(5.6), Inches(2.0), "Indian partners", [
        "Board / SOM qualification with SI & DPSU ecosystem.",
    ])

    # 14 Roadmap
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Roadmap")
    add_title(tf, "Phased engagement with LRDE")
    phases = [
        ("0", "Discovery — 4–6 weeks", "Workload survey (2 models), latency/SWaP envelopes, success metrics with ERP owners."),
        ("1", "PoC — 8–12 weeks", "Port models to Agilex eval kits; accuracy vs golden set; latency report for radar windows."),
        ("2", "Pilot — 3–6 months", "Rugged module or PCIe in processor cabinet; training-cluster slice; MLOps baseline."),
        ("3", "Scale — program-driven", "Production SOM/PCIe, environmental qualification, ILS & field update process."),
    ]
    for i, (num, title, body) in enumerate(phases):
        top = Inches(1.9) + i * Inches(1.25)
        circ = s.shapes.add_shape(9, MARGIN_X, top + Inches(0.15), Inches(0.55), Inches(0.55))  # oval
        fill_shape(circ, STEEL)
        nbox = add_textbox(s, MARGIN_X, top + Inches(0.25), Inches(0.55), Inches(0.4))
        p = nbox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = num
        set_run(p.runs[0], 14, True, WHITE)
        add_rect(s, MARGIN_X + Inches(0.75), top, content_w - Inches(0.75), Inches(1.1), PANEL)
        t = add_textbox(s, MARGIN_X + Inches(0.95), top + Inches(0.15), content_w - Inches(1.2), Inches(0.85))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        set_run(p.runs[0], 16, True, INK)
        add_body(tf, body, MUTED, 13)

    # 15 BOM
    s = blank(prs)
    set_slide_bg(s, PAPER)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Reference BOM")
    add_title(tf, "What LRDE would procure")
    card(s, MARGIN_X, Inches(2.0), Inches(5.8), Inches(3.5), "Training lab", [
        "• GPU nodes + IB fabric",
        "• Encrypted storage",
        "• Hardened MLOps stack",
    ])
    card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(3.5), "Inference & tools", [
        "• Agilex 7 PCIe + Agilex 5 kits",
        "• Quartus Prime + FPGA AI Suite",
        "• OpenVINO + enablement services",
    ])
    note = add_textbox(s, MARGIN_X, Inches(5.9), content_w, Inches(0.7))
    p = note.text_frame.paragraphs[0]
    p.text = "Detailed commercial annex after discovery workshop. See BOM_Reference.md."
    set_run(p.runs[0], 14, False, MUTED)

    # 16 Next steps
    s = blank(prs)
    set_slide_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), W, H, DARK_BG)
    box = add_textbox(s, MARGIN_X, MARGIN_Y, content_w, Inches(1.2))
    tf = box.text_frame
    add_eyebrow(tf, "Call to action", RGBColor(0xB8, 0xD9, 0xCC))
    add_title(tf, "Recommended next steps", WHITE, 36)
    steps = [
        ("1. Workshop", "2-day architecture session at LRDE — nominate radar AI & ERP stakeholders."),
        ("2. Select models", "Pick two flagship networks (e.g. pulse ID + RD classifier) for Phase-1 PoC."),
        ("3. Lock gates", "Agree accuracy, window latency, and power budgets before compile."),
        ("4. SOW", "Issue technical SOW + commercial annex for PoC kits and FAE support."),
    ]
    for i, (title, body) in enumerate(steps):
        col = i % 2
        row = i // 2
        left = MARGIN_X + col * Inches(6.2)
        top = Inches(2.0) + row * Inches(2.0)
        shape = add_rect(s, left, top, Inches(5.9), Inches(1.7), RGBColor(0x1A, 0x2E, 0x3A))
        t = add_textbox(s, left + Inches(0.25), top + Inches(0.25), Inches(5.4), Inches(1.3))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        set_run(p.runs[0], 16, True, WHITE)
        add_body(tf, body, RGBColor(0xC5, 0xD5, 0xCD), 13)
    meta = add_textbox(s, MARGIN_X, Inches(6.4), content_w, Inches(0.6))
    p = meta.text_frame.paragraphs[0]
    p.text = "Prepared for LRDE, DRDO Bengaluru    ·    Inference: Altera Agilex + FPGA AI Suite"
    set_run(p.runs[0], 12, False, RGBColor(0xB0, 0xC4, 0xBA))

    out = Path("/workspace/proposal/LRDE_AI_Slide_Deck.pptx")
    prs.save(out)
    art = Path("/opt/cursor/artifacts/LRDE_AI_Slide_Deck.pptx")
    art.parent.mkdir(parents=True, exist_ok=True)
    prs.save(art)
    print(f"Wrote {out} ({out.stat().st_size} bytes)")
    print(f"Wrote {art}")


if __name__ == "__main__":
    build()
